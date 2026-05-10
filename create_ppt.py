from pptx import Presentation
from pptx.util import Inches

prs = Presentation()

slides_data = [
    (
        "GEO: Generative Engine Optimization",
        "Analytics platform for AI Search Visibility"
    ),
    (
        "Problem Statement",
        "Traditional SEO tools cannot measure brand visibility across ChatGPT, Gemini, Claude, and Perplexity."
    ),
    (
        "Market Shift",
        "Search is evolving from link retrieval to AI-generated answers."
    ),
    (
        "Competitor Comparison",
        "Compare GEO with Ahrefs, Semrush, Profound, and Peec AI."
    ),
    (
        "Core Feature",
        "AI Narrative Drift Monitor tracks brand perception differences across LLMs."
    ),
    (
        "Technical Architecture",
        "Frontend: React\nBackend: Node.js\nAI APIs: OpenAI, Gemini\nDatabase: PostgreSQL"
    ),
    (
        "Monetization",
        "Freemium SaaS + Enterprise Analytics + API Access"
    ),
    (
        "Roadmap",
        "3 Months: AI visibility dashboard\n12 Months: autonomous GEO optimization"
    ),
    (
        "Conclusion",
        "SEO optimized websites. GEO optimizes machine understanding."
    )
]

for title, content in slides_data:
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)

    slide.shapes.title.text = title
    slide.placeholders[1].text = content

ppt_path = "GEO_Product_Strategy.pptx"
prs.save(ppt_path)

print(f"Presentation saved: {ppt_path}")